from pathlib import Path
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation
from langchain.schema import Document

def load_documents(data_dir: str) -> list[Document]:
    """
    Scans data_dir/raw for .pdf and .pptx files,
    loads them into Document objects, and returns a combined list.
    """
    docs: list[Document] = []
    raw_folder = Path(data_dir) / "raw"

    # ——— Load PDFs via PyPDFLoader ———
    for pdf_path in raw_folder.glob("*.pdf"):
        loader = PyPDFLoader(str(pdf_path))
        pdf_docs = loader.load()
        for page_num, doc in enumerate(pdf_docs, start=1):
            # annotate metadata
            doc.metadata.update({
                "source": pdf_path.name,
                "page": page_num,
            })
        docs.extend(pdf_docs)
        print(f"Loaded {len(pdf_docs)} pages from {pdf_path.name}")

    # ——— Load PPTX via python-pptx ———
    for ppt_path in raw_folder.glob("*.pptx"):
        prs = Presentation(str(ppt_path))
        ppt_docs: list[Document] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            content = "\n".join(texts)
            metadata = {
                "source": ppt_path.name,
                "slide_number": slide_idx,
            }
            ppt_docs.append(Document(page_content=content, metadata=metadata))
        docs.extend(ppt_docs)
        print(f"Loaded {len(ppt_docs)} slides from {ppt_path.name}")

    return docs

if __name__ == "__main__":
    docs = load_documents("data")
    print(f"Total documents loaded: {len(docs)}")
    